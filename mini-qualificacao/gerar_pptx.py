"""Gera a apresentacao de 8 slides (4:3, light mode, tipografia grande para projetor)
a partir dos mesmos insumos numericos usados pelo docx (notebooks 01-05).

Identidade visual UFABC: verde #006633 e amarelo #FFCC00 (extraidos dos logotipos
oficiais em `identidade visual_UFABC/img/`).

O layout e calculado: `altura_texto()` estima quantas linhas cada bloco ocupa e as
caixas/marcadores sao empilhados a partir dessa estimativa, em vez de alturas fixas
(que estouravam quando o texto quebrava em mais linhas do que o previsto).

Uso:  .venv/Scripts/python.exe mini-qualificacao/gerar_pptx.py
Saida: mini-qualificacao/INF-009 2026.2 - Apresentacao - Priscila Santos.pptx
       mini-qualificacao/insumos_gerados/figs_slides/*.png
"""
import math
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib.ticker import PercentFormatter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.metrics import confusion_matrix

ROOT = Path(__file__).resolve().parent.parent
MQ = ROOT / "mini-qualificacao"
INS = MQ / "insumos_gerados"
FIGS = INS / "figs_slides"
FEAT = ROOT / "datasets" / "DAiSEE" / "features"
LOGO = ROOT / "identidade visual_UFABC" / "img" / "logotipo-ufabc-lateral.png"
SAIDA = MQ / "INF-009 2026.2 - Apresentacao - Priscila Santos.pptx"

# --------------------------------------------------------------------------
# Tokens de identidade visual
# --------------------------------------------------------------------------
VERDE = "#006633"        # UFABC - cor primaria e slot categorico 1
AMARELO = "#FFCC00"      # UFABC - acento de chrome (reguas, realces), nunca marca sobre branco
OURO = "#A67C00"         # slot categorico 2: familia do amarelo, escurecido p/ contraste >= 3:1
VERDE_CLARO = "#E6F0EA"  # superficie de painel
VERDE_MEDIO = "#4C9970"  # degrau intermediario da rampa sequencial
CINZA = "#8A9490"        # de-enfase (contexto)
CINZA_CLARO = "#D7DCD9"  # grade / reguas
TINTA = "#14261C"        # texto principal
TINTA_FRACA = "#5B6660"  # texto secundario
BRANCO = "#FFFFFF"

FONTE = "Arial"

# Geometria do slide 4:3
SLIDE_W, SLIDE_H = 10.0, 7.5
MARGEM = 0.55
UTIL = SLIDE_W - 2 * MARGEM          # 8.90"
Y_CONTEUDO = 1.34                    # inicio do conteudo (sem subtitulo)
Y_CONTEUDO_SUB = 1.52                # inicio do conteudo (com subtitulo)
Y_RODAPE = 6.95


def rgb(hexa):
    return RGBColor.from_string(hexa.lstrip("#").upper())


# --------------------------------------------------------------------------
# Estimativa de altura de texto (Arial) - base de todo o layout
# --------------------------------------------------------------------------
def n_linhas(txt, size, larg, bold=False):
    """Numero estimado de linhas de `txt` em Arial `size` pt dentro de `larg` polegadas."""
    largura_char = size * (0.530 if bold else 0.505) / 72.0
    por_linha = max(int((larg / largura_char) * 0.94), 8)   # 6% de folga
    linhas = 0
    for paragrafo in str(txt).split("\n"):
        linhas += max(1, math.ceil(len(paragrafo) / por_linha))
    return linhas


def altura_texto(txt, size, larg, bold=False, entrelinhas=1.0):
    return n_linhas(txt, size, larg, bold) * size * 1.21 * entrelinhas / 72.0


def plano(txt):
    """Achata um conteudo rico (lista de (texto, bold)) em string."""
    if isinstance(txt, list):
        return "".join(t for t, _ in txt)
    return str(txt)


# --------------------------------------------------------------------------
# Carga dos insumos
# --------------------------------------------------------------------------
def carregar():
    d = {
        "consolidado": pd.read_csv(INS / "tabela_resultados_consolidados.csv"),
        "boot": pd.read_csv(INS / "tabela_bootstrap_ic95.csv"),
        "binario": pd.read_csv(INS / "tabela_objetivo2_binario.csv").iloc[0],
        "decisao": pd.read_csv(INS / "decisao_hipotese.csv").iloc[0],
        "importancia": pd.read_csv(INS / "tabela_importancia_features.csv", index_col=0),
        "estabilidade": pd.read_csv(INS / "tabela_estabilidade_features.csv", index_col=0),
        "nclip": pd.read_csv(FEAT / "n_clipes_por_split.csv").iloc[0],
        "taxa": pd.read_csv(FEAT / "taxa_deteccao_rosto.csv").iloc[0],
        "pred_a": pd.read_csv(FEAT / "predicoes_teste_trilha_a.csv"),
    }
    d["n_features"] = len(pd.read_csv(FEAT / "importancias_features_engagement.csv"))
    return d


def n(x, casas=3):
    return f"{x:.{casas}f}".replace(".", ",")


def pct(x, casas=1):
    return f"{100 * x:.{casas}f}".replace(".", ",") + "%"


def mil(x):
    return f"{int(x):,}".replace(",", ".")


# --------------------------------------------------------------------------
# Figuras (matplotlib) - tipografia grande, light mode, paleta UFABC validada
# --------------------------------------------------------------------------
def estilo_base():
    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.sans-serif": [FONTE, "DejaVu Sans"],
        "font.size": 15,
        "axes.titlesize": 16,
        "axes.labelsize": 15,
        "xtick.labelsize": 14,
        "ytick.labelsize": 15,
        "legend.fontsize": 14,
        "figure.facecolor": BRANCO,
        "axes.facecolor": BRANCO,
        "savefig.facecolor": BRANCO,
        "axes.edgecolor": CINZA_CLARO,
        "axes.linewidth": 0.8,
        "text.color": TINTA,
        "axes.labelcolor": TINTA_FRACA,
        "xtick.color": TINTA_FRACA,
        "ytick.color": TINTA_FRACA,
        "xtick.major.size": 0,
        "ytick.major.size": 0,
        "grid.color": CINZA_CLARO,
        "grid.linewidth": 0.8,
    })


def limpar(ax, eixos=("top", "right", "left")):
    for lado in eixos:
        ax.spines[lado].set_visible(False)


def salvar(fig, nome):
    FIGS.mkdir(parents=True, exist_ok=True)
    caminho = FIGS / nome
    fig.savefig(caminho, dpi=200, bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)
    return caminho


def fig_distribuicao():
    """Enfase: as duas classes raras sao o problema; as demais sao contexto."""
    rot = ["0\nMuito baixo", "1\nBaixo", "2\nAlto", "3\nMuito alto"]
    v = np.array([0.68, 5.10, 49.55, 44.67])
    cores = [OURO, OURO, CINZA, CINZA]
    fig, ax = plt.subplots(figsize=(5.35, 3.3))
    barras = ax.bar(rot, v, color=cores, width=0.62)
    for b, val in zip(barras, v):
        ax.text(b.get_x() + b.get_width() / 2, val + 1.6,
                f"{val:.2f}%".replace(".", ","), ha="center", va="bottom",
                fontsize=15, fontweight="bold", color=TINTA)
    ax.set_ylim(0, 60)
    ax.set_ylabel("% dos clipes")
    ax.yaxis.set_major_formatter(PercentFormatter())
    ax.set_yticks([0, 25, 50])
    ax.grid(axis="y")
    ax.set_axisbelow(True)
    limpar(ax)
    ax.set_title("Engajamento no DAiSEE (n = 8.925 clipes)", color=TINTA, pad=10)
    ax.annotate("Desengajado:\n5,8% do total", xy=(0.5, 8), xytext=(0.78, 30),
                fontsize=14, color=OURO, fontweight="bold", ha="center",
                arrowprops=dict(arrowstyle="-", color=OURO, lw=1.6))
    return salvar(fig, "fig_distribuicao.png")


def fig_sota():
    """Acuracia publicada por ano: literatura em cinza, nossa PoC em destaque."""
    pontos = [(2016, 57.9), (2021, 58.8), (2021, 67.4), (2022, 63.6), (2023, 66.6),
              (2024, 68.6), (2024, 70.2), (2025, 63.9), (2025, 73.4)]
    fig, ax = plt.subplots(figsize=(5.0, 3.5))
    ax.scatter([p[0] for p in pontos], [p[1] for p in pontos], s=110, color=CINZA, zorder=3,
               edgecolor=BRANCO, linewidth=2, label="Literatura (só acurácia)")
    ax.scatter([2026], [52.7], s=240, color=VERDE, zorder=4, marker="D",
               edgecolor=BRANCO, linewidth=2, label="PoC (acurácia + kappa + IC)")
    ax.annotate("ViBED-Net 73,4%", xy=(2025, 73.4), xytext=(-10, 12),
                textcoords="offset points", fontsize=13, color=TINTA_FRACA, ha="right")
    ax.annotate("LRCN 57,9%", xy=(2016, 57.9), xytext=(4, 12), textcoords="offset points",
                fontsize=13, color=TINTA_FRACA, ha="left")
    ax.annotate("PoC 52,7%", xy=(2026, 52.7), xytext=(0, 14), textcoords="offset points",
                fontsize=14, color=VERDE, fontweight="bold", ha="center")
    ax.set_ylim(45, 84)
    ax.set_xlim(2014.8, 2027.6)
    ax.set_xticks([2016, 2019, 2022, 2025])
    ax.set_ylabel("Acurácia relatada (%)")
    ax.grid(axis="y")
    ax.set_axisbelow(True)
    limpar(ax)
    ax.legend(loc="lower left", frameon=False, fontsize=13, handletextpad=0.2,
              borderaxespad=0.1)
    ax.set_title("DAiSEE · engajamento, 4 classes", color=TINTA, pad=10)
    return salvar(fig, "fig_sota.png")


NOMES_CURTOS = {
    "RandomForest (class_weight=balanced)": "Random Forest balanceada",
    "XGBoost (sample_weight balanceado)": "XGBoost balanceado",
    "RandomForest + SMOTE": "Random Forest + SMOTE",
    "BiLSTM ordinal (CORAL)": "BiLSTM ordinal (CORAL)",
    "RandomForest (sem balanceamento)": "Random Forest sem balanceio",
    "SVM (class_weight=balanced)": "SVM balanceada",
    "Baseline trivial (classe majoritaria)": "Baseline trivial",
    "BiLSTM + Focal Loss (class-weighted)": "BiLSTM + perda focal",
}


def fig_modelos(cons):
    """Duas metricas na mesma escala 0-1 -> um unico eixo, barras agrupadas."""
    df = cons.copy().iloc[::-1]
    rotulos = [NOMES_CURTOS.get(m.strip(), m) for m in df["modelo"]]
    y = np.arange(len(df))
    alt = 0.36
    fig, ax = plt.subplots(figsize=(5.5, 3.8))
    ax.barh(y + alt / 2 + 0.02, df["macro_f1"], height=alt, color=VERDE, label="macro-F1")
    ax.barh(y - alt / 2 - 0.02, df["kappa"], height=alt, color=OURO, label="kappa")
    for yy, v in zip(y + alt / 2 + 0.02, df["macro_f1"]):
        ax.text(v + 0.008, yy, n(v, 2), va="center", fontsize=13, color=TINTA)
    ax.set_yticks(y)
    ax.set_yticklabels(rotulos, fontsize=13.5)
    for rot, nome in zip(ax.get_yticklabels(), rotulos):
        if nome == "Baseline trivial":
            rot.set_color(TINTA_FRACA)
            rot.set_style("italic")
    ax.set_xlim(0, 0.46)
    ax.set_xlabel("Valor da métrica (conjunto de teste)")
    ax.grid(axis="x")
    ax.set_axisbelow(True)
    limpar(ax, ("top", "right"))
    ax.spines["left"].set_color(CINZA_CLARO)
    ax.legend(loc="lower right", frameon=False, fontsize=13.5)
    return salvar(fig, "fig_modelos.png")


def fig_bootstrap(boot):
    """IC 95% por bootstrap: o teste da hipotese e a nao sobreposicao."""
    fig, ax = plt.subplots(figsize=(4.3, 3.5))
    linhas = [("Melhor modelo", boot.iloc[0], VERDE), ("Baseline trivial", boot.iloc[1], CINZA)]
    for i, (nome, r, cor) in enumerate(linhas):
        y = 1 - i
        ax.plot([r["macro_f1_ic_low"], r["macro_f1_ic_high"]], [y, y], color=cor, lw=5,
                solid_capstyle="round", zorder=2)
        ax.plot([r["macro_f1_media"]], [y], "o", color=cor, ms=13, markeredgecolor=BRANCO,
                markeredgewidth=2, zorder=3)
        ax.text(r["macro_f1_media"], y + 0.22, n(r["macro_f1_media"], 3), ha="center",
                fontsize=15, fontweight="bold", color=cor)
        ax.text(0.015, y - 0.34, nome, fontsize=14, color=TINTA_FRACA)
        ax.text(0.015, y - 0.60, f"IC 95% [{n(r['macro_f1_ic_low'], 3)}; "
                                 f"{n(r['macro_f1_ic_high'], 3)}]",
                fontsize=12.5, color=CINZA)
    ax.axvspan(boot.iloc[1]["macro_f1_ic_high"], boot.iloc[0]["macro_f1_ic_low"],
               color=VERDE_CLARO, zorder=0)
    ax.text((boot.iloc[1]["macro_f1_ic_high"] + boot.iloc[0]["macro_f1_ic_low"]) / 2, -0.98,
            "sem sobreposição", ha="center", fontsize=13.5, color=VERDE, fontweight="bold")
    ax.set_ylim(-1.25, 1.55)
    ax.set_xlim(0.0, 0.55)
    ax.set_yticks([])
    ax.set_xlabel("macro-F1 · IC 95% (bootstrap)", fontsize=13.5)
    ax.grid(axis="x")
    ax.set_axisbelow(True)
    limpar(ax, ("top", "right", "left"))
    return salvar(fig, "fig_bootstrap.png")


ROTULO_FEAT = {
    "ear_mean": "EAR médio (abertura ocular)",
    "ear_min": "EAR mínimo",
    "gaze_v_std": "Variação do olhar (vertical)",
    "ear_std": "Variação do EAR",
    "gaze_h_std": "Variação do olhar (horizontal)",
    "ear_max": "EAR máximo",
    "roll_std": "Variação da inclinação (roll)",
    "gaze_h_min": "Olhar horizontal mínimo",
}


def fig_features(imp, est):
    """Enfase: atributos estaveis (>= 4/5 folds) em verde, instaveis em cinza."""
    top = imp.head(8).iloc[::-1]
    freq = est.iloc[:, 0]
    cores = [VERDE if freq.get(i, 0) >= 4 else CINZA for i in top.index]
    rotulos = [ROTULO_FEAT.get(i, i) for i in top.index]
    fig, ax = plt.subplots(figsize=(5.4, 3.5))
    barras = ax.barh(np.arange(len(top)), top["importancia"], color=cores, height=0.62)
    for b, v in zip(barras, top["importancia"]):
        ax.text(v + 0.0015, b.get_y() + b.get_height() / 2, n(v, 3), va="center",
                fontsize=13, color=TINTA)
    ax.set_yticks(np.arange(len(top)))
    ax.set_yticklabels(rotulos, fontsize=13.5)
    ax.set_xlim(0, 0.105)
    ax.set_xticks([0, 0.04, 0.08])
    ax.set_xlabel("Importância (Random Forest)")
    ax.grid(axis="x")
    ax.set_axisbelow(True)
    limpar(ax, ("top", "right"))
    ax.spines["left"].set_color(CINZA_CLARO)
    ax.set_title("Em verde: estável no top-5 em ≥ 4 de 5 folds", color=VERDE, fontsize=13.5,
                 loc="left", pad=8)
    return salvar(fig, "fig_features.png")


def fig_confusao(pred):
    """Matriz binaria normalizada por linha: rampa sequencial de um unico matiz."""
    yb_true = (pred["y_true"] >= 2).astype(int)
    yb_pred = (pred["y_pred"] >= 2).astype(int)
    cm = confusion_matrix(yb_true, yb_pred)
    cmn = cm / cm.sum(axis=1, keepdims=True)
    rampa = matplotlib.colors.LinearSegmentedColormap.from_list(
        "ufabc", [BRANCO, VERDE_CLARO, VERDE_MEDIO, VERDE])
    fig, ax = plt.subplots(figsize=(4.3, 3.4))
    ax.imshow(cmn, cmap=rampa, vmin=0, vmax=1)
    rot = ["Desengajado", "Engajado"]
    for i in range(2):
        for j in range(2):
            claro = cmn[i, j] > 0.55
            ax.text(j, i - 0.10, pct(cmn[i, j], 1), ha="center", va="center", fontsize=19,
                    fontweight="bold", color=BRANCO if claro else TINTA)
            ax.text(j, i + 0.20, f"{cm[i, j]} clipes", ha="center", va="center", fontsize=13,
                    color=BRANCO if claro else TINTA_FRACA)
    ax.set_xticks([0, 1], rot, fontsize=14)
    ax.set_yticks([0, 1], rot, fontsize=14, rotation=90, va="center")
    ax.set_xlabel("Previsto", fontsize=14)
    ax.set_ylabel("Real", fontsize=14)
    for lado in ("top", "right", "bottom", "left"):
        ax.spines[lado].set_visible(False)
    return salvar(fig, "fig_confusao_binaria.png")


# --------------------------------------------------------------------------
# Infraestrutura de slides
# --------------------------------------------------------------------------
def texto(slide, x, y, cx, cy, blocos, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          entrelinhas=0.98):
    """blocos = lista de dicts {t, size, bold, color, space_after, entrelinhas}.
    `t` pode ser str ou lista de (texto, bold) para negrito seletivo."""
    caixa = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cx), Inches(cy))
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(blocos):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = b.get("entrelinhas", entrelinhas)
        p.space_after = Pt(b.get("space_after", 5))
        p.space_before = Pt(b.get("space_before", 0))
        conteudo = b["t"]
        partes = conteudo if isinstance(conteudo, list) else [(conteudo, b.get("bold", False))]
        for txt, negrito in partes:
            r = p.add_run()
            r.text = txt
            r.font.name = FONTE
            r.font.size = Pt(b.get("size", 18))
            r.font.bold = negrito
            r.font.color.rgb = rgb(b.get("color", TINTA))
    return caixa


def altura_blocos(blocos, larg, pad=0.0):
    """Altura total estimada de uma lista de blocos de `texto()`."""
    total = pad
    for b in blocos:
        size = b.get("size", 18)
        negrito = b.get("bold", False) or (
            isinstance(b["t"], list) and any(n_ for _, n_ in b["t"]))
        total += altura_texto(plano(b["t"]), size, larg, negrito,
                              b.get("entrelinhas", 1.0))
        total += b.get("space_after", 5) / 72.0
    return total


def retangulo(slide, x, y, cx, cy, cor, linha=None, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cx), Inches(cy))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(cor)
    if linha:
        s.line.color.rgb = rgb(linha)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def novo_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # layout em branco


def cabecalho(slide, titulo, sub=None):
    """Titulo de uma linha (27 pt). Avisa no console se a estimativa passar de 1 linha."""
    if n_linhas(titulo, 27, UTIL, bold=True) > 1:
        print(f"  [aviso] titulo longo demais para uma linha: {titulo!r}")
    texto(slide, MARGEM, 0.34, UTIL, 0.55,
          [{"t": titulo, "size": 27, "bold": True, "color": VERDE, "space_after": 0}])
    retangulo(slide, MARGEM, 0.95, 1.5, 0.055, AMARELO)
    retangulo(slide, MARGEM + 1.62, 0.973, UTIL - 1.62, 0.011, CINZA_CLARO)
    if sub:
        texto(slide, MARGEM, 1.10, UTIL, 0.32,
              [{"t": sub, "size": 15.5, "color": TINTA_FRACA, "space_after": 0}])


def rodape(slide, numero):
    retangulo(slide, MARGEM, Y_RODAPE, UTIL, 0.011, CINZA_CLARO)
    texto(slide, MARGEM, Y_RODAPE + 0.11, 7.2, 0.3,
          [{"t": "Classificação de Engajamento em EaD com Visão Computacional · INF-009 · UFABC",
            "size": 10.5, "color": TINTA_FRACA, "space_after": 0}])
    texto(slide, SLIDE_W - MARGEM - 1.0, Y_RODAPE + 0.11, 1.0, 0.3,
          [{"t": str(numero), "size": 10.5, "bold": True, "color": VERDE, "space_after": 0}],
          align=PP_ALIGN.RIGHT)


def painel(slide, x, y, cx, blocos, fundo=VERDE_CLARO, borda=None, faixa=None,
           pad_x=0.26, pad_y=0.16, altura=None):
    """Caixa com altura calculada a partir do conteudo. Retorna a altura usada."""
    larg_txt = cx - 2 * pad_x
    cy = altura if altura else altura_blocos(blocos, larg_txt) + 2 * pad_y
    retangulo(slide, x, y, cx, cy, fundo, linha=borda, lw=1.75)
    if faixa:
        retangulo(slide, x, y, 0.07, cy, faixa)
    texto(slide, x + pad_x, y + pad_y, larg_txt, cy - 2 * pad_y, blocos)
    return cy


def marcadores(slide, x, y, cx, itens, size=18, cor_marca=VERDE, gap=0.14):
    """Empilha itens com marcador quadrado, cada um com a altura que realmente ocupa."""
    larg_txt = cx - 0.30
    yy = y
    for item in itens:
        alt = altura_texto(plano(item), size, larg_txt, bold=True, entrelinhas=1.0)
        retangulo(slide, x, yy + size * 0.006, 0.13, 0.13, cor_marca)
        texto(slide, x + 0.30, yy, larg_txt, alt,
              [{"t": item if isinstance(item, list) else [(item, False)], "size": size,
                "color": TINTA, "space_after": 0, "entrelinhas": 1.0}])
        yy += alt + gap
    return yy - y


def estatistica(slide, x, y, cx, valor, rotulo, size_valor=24):
    texto(slide, x, y, cx, 0.45,
          [{"t": valor, "size": size_valor, "bold": True, "color": VERDE, "space_after": 0,
            "entrelinhas": 0.9}], align=PP_ALIGN.CENTER)
    texto(slide, x, y + 0.42, cx, 0.5,
          [{"t": rotulo, "size": 13, "color": TINTA_FRACA, "space_after": 0,
            "entrelinhas": 0.98}], align=PP_ALIGN.CENTER)


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------
def slide_1_capa(prs):
    s = novo_slide(prs)
    retangulo(s, 0, 0, SLIDE_W, 0.30, VERDE)
    retangulo(s, 0, 0.30, SLIDE_W, 0.09, AMARELO)
    s.shapes.add_picture(str(LOGO), Inches(MARGEM), Inches(0.72), width=Inches(2.75))

    texto(s, MARGEM, 2.55, 8.9, 2.0, [
        {"t": "Classificação de Engajamento", "size": 40, "bold": True, "color": VERDE,
         "space_after": 0, "entrelinhas": 0.92},
        {"t": "em EaD com Visão Computacional", "size": 40, "bold": True, "color": VERDE,
         "space_after": 0, "entrelinhas": 0.92},
    ])
    retangulo(s, MARGEM, 4.38, 2.2, 0.07, AMARELO)
    texto(s, MARGEM, 4.62, 8.6, 0.9, [
        {"t": "Prova de conceito no DAiSEE: marcos faciais leves, sem transmitir imagem "
              "do rosto, sob protocolo multimétrica", "size": 18, "color": TINTA_FRACA,
         "space_after": 0, "entrelinhas": 1.05},
    ])
    texto(s, MARGEM, 5.62, 8.9, 1.3, [
        {"t": "Priscila Henrique Medeiro dos Santos", "size": 21, "bold": True,
         "color": TINTA, "space_after": 4},
        {"t": "INF-009 – Projeto e Comunicação de Pesquisa em Engenharia da Informação",
         "size": 15, "color": TINTA_FRACA, "space_after": 2},
        {"t": "Pós-graduação em Engenharia da Informação · UFABC · Santo André, agosto de 2026",
         "size": 15, "color": TINTA_FRACA, "space_after": 0},
    ])
    retangulo(s, 0, 7.28, SLIDE_W, 0.22, VERDE)
    return s


def slide_2_problema(prs, figs):
    s = novo_slide(prs)
    cabecalho(s, "O problema: desengajamento invisível",
              "Por que a métrica dominante na literatura esconde justamente quem precisa de atenção")
    marcadores(s, MARGEM, Y_CONTEUDO_SUB + 0.06, 4.15, [
        [("Em aula síncrona, o docente ", False),
         ("não observa dezenas de participantes ao mesmo tempo", True), (".", False)],
        [("A queda de atenção é percebida tarde — e alimenta a ", False),
         ("evasão em EaD", True), (".", False)],
        [("No DAiSEE, ", False), ("94% dos clipes são “engajado”", True),
         (": responder sempre “engajado” já acerta 51,8% de acurácia.", False)],
    ], size=18, gap=0.24)

    s.shapes.add_picture(str(figs["dist"]), Inches(4.78), Inches(Y_CONTEUDO_SUB),
                         width=Inches(4.72))

    painel(s, MARGEM, 5.22, UTIL, [
        {"t": [("Hipótese.  ", True),
               ("Indicadores visuais analisados continuamente — abertura ocular, direção do "
                "olhar, pose da cabeça — identificam níveis de engajamento com mais "
                "objetividade que o acompanhamento humano não instrumentado.", False)],
         "size": 17.5, "color": TINTA, "space_after": 0, "entrelinhas": 1.05},
    ], fundo=BRANCO, borda=VERDE, pad_y=0.20)
    rodape(s, 2)
    return s


def slide_3_objetivos(prs):
    s = novo_slide(prs)
    cabecalho(s, "Objetivos mensuráveis, definidos a priori",
              "Cada objetivo traz o critério numérico que o declara atingido — ou não")
    cartoes = [
        ("Objetivo geral",
         "Avaliar se indicadores visuais permitem identificar níveis de engajamento em aulas "
         "síncronas de EaD.",
         "superar o baseline de classe majoritária em macro-F1 e kappa, com IC 95% "
         "(bootstrap, 1.000 reamostragens) não sobrepostos."),
        ("Objetivo específico 1",
         "Identificar os indicadores visuais mais relacionados ao engajamento, ao "
         "desengajamento e à distração.",
         "ranking de importância com top-5 estável em ≥ 4 de 5 folds de validação cruzada."),
        ("Objetivo específico 2",
         "Diferenciar engajados (níveis 2–3) de desengajados (níveis 0–1) em vídeo real.",
         "revocação da classe Desengajado ≥ 60% e AUC-ROC ≥ 0,70, contra revocação próxima "
         "de zero do baseline trivial."),
    ]
    y = Y_CONTEUDO_SUB + 0.02
    for titulo, corpo, criterio in cartoes:
        alt = painel(s, MARGEM, y, UTIL, [
            {"t": titulo, "size": 18, "bold": True, "color": VERDE, "space_after": 3},
            {"t": corpo, "size": 15.5, "color": TINTA, "space_after": 2, "entrelinhas": 1.0},
            {"t": [("Critério: ", True), (criterio, False)],
             "size": 15.5, "color": TINTA, "space_after": 0, "entrelinhas": 1.0},
        ], faixa=VERDE, pad_y=0.13)
        y += alt + 0.13
    rodape(s, 3)
    return s


def slide_4_sota(prs, figs):
    s = novo_slide(prs)
    cabecalho(s, "Estado da arte 2016 → 2026",
              "Muito modelo, pouca métrica: onde a literatura do DAiSEE ainda não olha")
    s.shapes.add_picture(str(figs["sota"]), Inches(0.40), Inches(Y_CONTEUDO_SUB + 0.12),
                         width=Inches(4.85))
    x = 5.52
    texto(s, x, Y_CONTEUDO_SUB + 0.02, 3.95, 0.35,
          [{"t": "As lacunas que atacamos", "size": 18, "bold": True, "color": VERDE,
            "space_after": 0}])
    marcadores(s, x, Y_CONTEUDO_SUB + 0.44, 3.95, [
        [("11 de 11", True), (" trabalhos reportam acurácia; ", False), ("nenhum", True),
         (" reporta kappa ou IC.", False)],
        [("Desbalanceamento extremo tratado implicitamente.", False)],
        [("Ordinalidade do rótulo (0→3) pouco explorada.", False)],
        [("SOTA (73,4%) exige treino pesado, fora do alcance de uma PoC.", False)],
        [("Zero-shot com VLMs falha: ", False), ("kappa < 0,10", True), (".", False)],
    ], size=15.5, gap=0.11)
    painel(s, 0.40, 5.68, 9.10, [
        {"t": [("Nossa aposta.  ", True),
               ("Não superar o número absoluto do SOTA, e sim medir honestamente o que "
                "atributos faciais leves entregam — sob um protocolo que torne visível o "
                "efeito do desbalanceamento.", False)],
         "size": 16.5, "color": BRANCO, "space_after": 0, "entrelinhas": 1.05},
    ], fundo=VERDE, pad_y=0.18)
    rodape(s, 4)
    return s


def slide_5_metodo(prs, meta):
    s = novo_slide(prs)
    cabecalho(s, "Metodologia: pipeline leve, sem imagem")
    x0, larg_util = 0.42, 9.16

    etapas = [("Vídeos DAiSEE", "clipes de 10 s\n640×480 · 30 fps"),
              ("20 quadros/clipe", "amostragem uniforme\n(não todos os quadros)"),
              ("MediaPipe FaceMesh + Íris", "EAR · MAR · pose\n(yaw, pitch, roll) · gaze")]
    y = Y_CONTEUDO + 0.06
    larg, gap, alt = 2.85, 0.30, 1.15
    for i, (t, sub) in enumerate(etapas):
        xx = x0 + i * (larg + gap)
        retangulo(s, xx, y, larg, alt, VERDE_CLARO)
        retangulo(s, xx, y, larg, 0.06, VERDE)
        texto(s, xx + 0.14, y + 0.20, larg - 0.28, alt - 0.26, [
            {"t": t, "size": 16, "bold": True, "color": VERDE, "space_after": 3,
             "entrelinhas": 0.95},
            {"t": sub, "size": 13.5, "color": TINTA_FRACA, "space_after": 0,
             "entrelinhas": 0.95}], align=PP_ALIGN.CENTER)
        if i < 2:
            texto(s, xx + larg, y + 0.34, gap, 0.4,
                  [{"t": "▶", "size": 15, "bold": True, "color": AMARELO, "space_after": 0}],
                  align=PP_ALIGN.CENTER)

    y2 = y + alt + 0.28
    trilhas = [
        (VERDE, "Trilha A · ML clássico",
         [f"{meta['n_features']} atributos agregados por clipe",
          "Random Forest · XGBoost · SVM",
          "GridSearchCV + 4 tratamentos de desbalanceamento"]),
        (OURO, "Trilha B · temporal raso",
         ["Sequência de 20 passos por clipe",
          "BiLSTM com perda focal ponderada",
          "Variante ordinal CORAL"]),
    ]
    larg2 = 4.43
    blocos_por_trilha = []
    for cor, t, linhas in trilhas:
        blocos = [{"t": t, "size": 16.5, "bold": True, "color": cor, "space_after": 5}]
        blocos += [{"t": l, "size": 14, "color": TINTA, "space_after": 2, "entrelinhas": 1.0}
                   for l in linhas]
        blocos_por_trilha.append((cor, blocos))
    alt2 = max(altura_blocos(b, larg2 - 0.4) for _, b in blocos_por_trilha) + 0.32
    for i, (cor, blocos) in enumerate(blocos_por_trilha):
        xx = x0 + i * (larg2 + 0.30)
        retangulo(s, xx, y2, larg2, alt2, BRANCO, linha=cor, lw=1.75)
        texto(s, xx + 0.20, y2 + 0.16, larg2 - 0.4, alt2 - 0.3, blocos)

    y3 = y2 + alt2 + 0.24
    retangulo(s, x0, y3, larg_util, 0.60, VERDE)
    texto(s, x0, y3, larg_util, 0.60, [
        {"t": "Avaliação: acurácia · macro-F1 · kappa · revocação por classe · IC 95% (bootstrap)",
         "size": 16, "bold": True, "color": BRANCO, "space_after": 0}],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    y4 = y3 + 0.60 + 0.26
    stats = [(mil(meta["n_total"]), "clipes processados", 24),
             ("69 / 19 / 21", "sujeitos disjuntos\ntreino / val. / teste", 21),
             (pct(meta["taxa_deteccao_media"], 2), "detecção facial", 24),
             ("42", "semente fixa\n(reprodutibilidade)", 24)]
    for i, (v, r, sz) in enumerate(stats):
        estatistica(s, x0 + i * (larg_util / 4), y4, larg_util / 4, v, r, size_valor=sz)
    rodape(s, 5)
    return s


def slide_6_resultados(prs, figs, cons, boot):
    s = novo_slide(prs)
    melhor = cons.iloc[0]
    base = cons[cons["modelo"].str.contains("trivial", case=False)].iloc[0]
    b_mod = boot.iloc[0]
    cabecalho(s, "Resultados: o sinal existe e é real",
              "Oito modelos no mesmo protocolo, comparados ao baseline de classe majoritária")
    s.shapes.add_picture(str(figs["modelos"]), Inches(0.30), Inches(Y_CONTEUDO_SUB - 0.04),
                         width=Inches(5.10))
    s.shapes.add_picture(str(figs["boot"]), Inches(5.78), Inches(Y_CONTEUDO_SUB + 0.06),
                         width=Inches(3.62))

    painel(s, 0.40, 5.30, 9.10, [
        {"t": [("Random Forest com class_weight balanceado: macro-F1 ", True),
               (n(melhor["macro_f1"], 3), True), (" e kappa ", True),
               (n(melhor["kappa"], 3), True),
               (f", contra {n(base['macro_f1'], 3)} e 0,000 do baseline. ", False),
               ("Os IC 95% não se sobrepõem — o ganho é significativo. Na Trilha B, a "
                "formulação ordinal CORAL superou a perda focal (0,266 contra 0,121).",
                False)],
         "size": 15, "color": TINTA, "space_after": 0, "entrelinhas": 1.05},
    ], fundo=BRANCO, borda=VERDE, pad_y=0.16)
    rodape(s, 6)
    return s


def slide_7_objetivos_12(prs, figs, binr, est):
    s = novo_slide(prs)
    cabecalho(s, "Objetivos 1 e 2: o que o rosto entrega",
              "Quais indicadores importam — e até onde eles separam engajado de desengajado")
    texto(s, 0.42, Y_CONTEUDO_SUB, 4.9, 0.35, [
        {"t": "OE1 ✓  atingido", "size": 17.5, "bold": True, "color": VERDE,
         "space_after": 0}])
    s.shapes.add_picture(str(figs["feat"]), Inches(0.32), Inches(Y_CONTEUDO_SUB + 0.36),
                         width=Inches(5.05))

    texto(s, 5.72, Y_CONTEUDO_SUB, 3.9, 0.35, [
        {"t": "OE2 ✗  parcialmente atingido", "size": 17.5, "bold": True, "color": OURO,
         "space_after": 0}])
    s.shapes.add_picture(str(figs["conf"]), Inches(5.98), Inches(Y_CONTEUDO_SUB + 0.42),
                         width=Inches(3.30))

    n_est = int((est.iloc[:, 0] >= 4).sum())
    esq = [
        {"t": f"{n_est} atributos no top-5 em ≥ 4 de 5 folds", "size": 15.5, "bold": True,
         "color": VERDE, "space_after": 3},
        {"t": "Abertura ocular (EAR) e variação do olhar dominam — sinais de sonolência e "
              "de dispersão do foco.", "size": 14, "color": TINTA, "space_after": 0,
         "entrelinhas": 1.02},
    ]
    dir_ = [
        {"t": [(f"AUC-ROC {n(binr['auc_roc'], 3)} ✓", True), ("  (meta ≥ 0,70)", False)],
         "size": 15.5, "color": VERDE, "space_after": 2},
        {"t": [(f"Revocação Desengajado {pct(binr['recall_desengajado'], 1)} ✗", True),
               ("  (meta ≥ 60%)", False)], "size": 15.5, "color": OURO, "space_after": 3},
        {"t": "O limiar padrão ainda perde 3 de cada 4 alunos desengajados.", "size": 14,
         "color": TINTA, "space_after": 0, "entrelinhas": 1.02},
    ]
    alt = max(altura_blocos(esq, 4.60 - 0.52), altura_blocos(dir_, 4.30 - 0.52)) + 0.30
    painel(s, 0.42, 5.22, 4.60, esq, fundo=VERDE_CLARO, altura=alt)
    painel(s, 5.20, 5.22, 4.30, dir_, fundo=BRANCO, borda=CINZA_CLARO, altura=alt)
    rodape(s, 7)
    return s


def slide_8_conclusao(prs, meta):
    s = novo_slide(prs)
    cabecalho(s, "Hipótese parcialmente corroborada")
    y_box = Y_CONTEUDO + 0.02
    alt_box = painel(s, 0.42, y_box, 9.10, [
        {"t": [("Veredito.  ", True),
               ("Há ganho estatisticamente significativo sobre o baseline em macro-F1 e "
                "kappa, mas a revocação da classe Desengajado ficou muito abaixo da meta.",
                False)],
         "size": 16.5, "color": BRANCO, "space_after": 0, "entrelinhas": 1.03},
    ], fundo=VERDE, pad_y=0.16, altura=0.94)

    y = y_box + alt_box + 0.20
    texto(s, 0.42, y, 4.5, 0.35, [
        {"t": "O que aprendemos", "size": 17.5, "bold": True, "color": VERDE,
         "space_after": 0}])
    marcadores(s, 0.42, y + 0.40, 4.45, [
        [("Marcos faciais leves ", False), ("carregam sinal real", True),
         (" de engajamento.", False)],
        [("kappa 0,110", True), (": concordância pouco acima do acaso — insuficiente para "
                                 "uso pedagógico direto.", False)],
        [("O gargalo é o desbalanceamento, não a arquitetura", True),
         (": nenhuma mitigação recuperou as classes raras.", False)],
    ], size=15, gap=0.13)

    texto(s, 5.20, y, 4.3, 0.35, [
        {"t": "Próximos passos (mestrado)", "size": 17.5, "bold": True, "color": VERDE,
         "space_after": 0}])
    marcadores(s, 5.20, y + 0.40, 4.30, [
        [("Meses 1–3: ", True), ("generalização entre bases (EngageNet).", False)],
        [("Meses 4–7: ", True), ("contexto de cena + aprendizado contrastivo ordinal.", False)],
        [("Meses 8–10: ", True), ("ruído de rótulo por sujeito e calibração de incerteza.",
                                  False)],
        [("Meses 11–18: ", True), ("painel de feedback ao docente e defesa.", False)],
    ], size=15, gap=0.13, cor_marca=OURO)

    painel(s, 0.42, 5.45, 9.10, [
        {"t": [("Limitações declaradas.  ", True),
               (f"O teste contém apenas {meta['n_classe0_teste']} clipes da classe 0, o que "
                "amplia os intervalos de confiança; 20 quadros por clipe descartam "
                "microdinâmicas; marcos faciais ignoram cena e postura corporal; há indícios "
                "de ruído de rótulo por sujeito.", False)],
         "size": 14.5, "color": TINTA, "space_after": 0, "entrelinhas": 1.03},
    ], fundo=BRANCO, borda=CINZA_CLARO, pad_y=0.14)
    rodape(s, 8)
    return s


# --------------------------------------------------------------------------
def main():
    d = carregar()
    estilo_base()

    nclip, taxa = d["nclip"], d["taxa"]
    total = int(nclip["train"] + nclip["validation"] + nclip["test"])
    meta = {
        "n_total": total,
        "n_test": int(nclip["test"]),
        "n_features": d["n_features"],
        "taxa_deteccao_media": float(
            (taxa["train"] * nclip["train"] + taxa["validation"] * nclip["validation"]
             + taxa["test"] * nclip["test"]) / total),
        "n_classe0_teste": int((d["pred_a"]["y_true"] == 0).sum()),
    }

    figs = {
        "dist": fig_distribuicao(),
        "sota": fig_sota(),
        "modelos": fig_modelos(d["consolidado"]),
        "boot": fig_bootstrap(d["boot"]),
        "feat": fig_features(d["importancia"], d["estabilidade"]),
        "conf": fig_confusao(d["pred_a"]),
    }

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)     # 4:3
    prs.slide_height = Inches(SLIDE_H)

    slide_1_capa(prs)
    slide_2_problema(prs, figs)
    slide_3_objetivos(prs)
    slide_4_sota(prs, figs)
    slide_5_metodo(prs, meta)
    slide_6_resultados(prs, figs, d["consolidado"], d["boot"])
    slide_7_objetivos_12(prs, figs, d["binario"], d["estabilidade"])
    slide_8_conclusao(prs, meta)

    prs.save(str(SAIDA))
    print("Apresentacao gerada:", SAIDA)
    print("Figuras em:", FIGS)
    return SAIDA


if __name__ == "__main__":
    main()
